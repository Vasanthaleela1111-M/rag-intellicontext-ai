from langchain_community.document_loaders import (
    PyMuPDFLoader,
    TextLoader,
    CSVLoader
)

from pptx import Presentation

from langchain_core.documents import Document
from docx import Document as DocxDocument

from utils.text_cleaner import clean_text


def load_docx(file_path):

    doc = DocxDocument(file_path)

    text = "\n".join(
        [para.text for para in doc.paragraphs]
    )

    return [
        Document(
            page_content=text,
            metadata={"source": file_path}
        )
    ]
def load_pptx(file_path):

    prs = Presentation(file_path)

    text = []

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text.append(shape.text)

    return [
        Document(
            page_content="\n".join(text),
            metadata={"source": file_path}
        )
    ]

def load_documents(file_path):

    if file_path.endswith(".pdf"):

        docs = PyMuPDFLoader(
            file_path
        ).load()

        for doc in docs:

            doc.page_content = clean_text(
                doc.page_content
            )

        return docs

    elif file_path.endswith(".txt"):
        return TextLoader(file_path).load()

    elif file_path.endswith(".csv"):
        return CSVLoader(file_path).load()

    elif file_path.endswith(".docx"):
        return load_docx(file_path)
    
    elif file_path.endswith(".pptx"):

        return load_pptx(file_path)

    else:
        raise ValueError(
            "Unsupported File Format"
        )
